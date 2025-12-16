import os
import json
import uuid
import time
from dataclasses import dataclass
from typing import Dict, Tuple, Any, List

from dotenv import load_dotenv
from openai import OpenAI

from pptx import Presentation

from telegram import (
    Update,
    InlineKeyboardButton,
    InlineKeyboardMarkup,
)
from telegram.ext import (
    Application,
    CommandHandler,
    MessageHandler,
    CallbackQueryHandler,
    PollAnswerHandler,
    ContextTypes,
    filters,
)

# ===================== CONFIG =====================
load_dotenv()

TELEGRAM_TOKEN = os.getenv("TELEGRAM_BOT_TOKEN")
OPENAI_KEY = os.getenv("OPENAI_API_KEY")

if not TELEGRAM_TOKEN:
    raise RuntimeError("TELEGRAM_BOT_TOKEN .env faylda yo'q")
if not OPENAI_KEY:
    raise RuntimeError("OPENAI_API_KEY .env faylda yo'q")

client = OpenAI(api_key=OPENAI_KEY)

QUESTIONS_PER_TEST = 25
TEST_DURATION_SECONDS = 25 * 60  # 25 daqiqa

SYSTEM_RULES = (
    "You are an academic study assistant in Uzbek. "
    "Create learning-oriented presentation content and tests. "
    "Avoid producing ready-to-submit cheating essays; instead give educational structure."
)

# Session key: (chat_id, user_id)
SessionKey = Tuple[int, int]


def _now() -> float:
    return time.time()


def get_sessions(app: Application) -> Dict[SessionKey, Dict[str, Any]]:
    # Saqlash joyi: application.bot_data
    if "sessions" not in app.bot_data:
        app.bot_data["sessions"] = {}
    return app.bot_data["sessions"]


# ===================== AI HELPERS =====================
def chat(prompt: str, temperature: float = 0.4) -> str:
    resp = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "system", "content": SYSTEM_RULES},
            {"role": "user", "content": prompt},
        ],
        temperature=temperature,
    )
    return resp.choices[0].message.content.strip()


def _strip_json_fence(text: str) -> str:
    t = text.strip()
    t = t.replace("```json", "```").strip()
    if t.startswith("```") and t.endswith("```"):
        t = t[3:-3].strip()
    return t.strip()


def generate_slide_plan(topic: str) -> dict:
    prompt = f"""
Mavzu: {topic}

10-12 slayddan iborat prezentatsiya rejasini JSON formatda qaytar:
- title: umumiy prezentatsiya sarlavhasi
- slides: ro'yxat
  - title: slayd sarlavhasi
  - bullets: 3-6 ta qisqa punkt
  - notes: 1-2 gaplik spiker eslatma

FAKAT JSON qaytar. Hech qanday izoh yozma.
"""
    raw = chat(prompt, temperature=0.3)
    raw = _strip_json_fence(raw)
    return json.loads(raw)


def generate_test_questions(topic: str) -> List[dict]:
    prompt = f"""
Mavzu: {topic}

{QUESTIONS_PER_TEST} ta test savoli tuz.
Har savol 4 variantli bo'lsin (A,B,C,D) va to'g'ri javob ham berilsin.
Natijani JSON array ko'rinishida qaytar.

Format:
[
  {{"q":"...","options":["A) ...","B) ...","C) ...","D) ..."],"answer":"A"}},
  ...
]

FAKAT JSON qaytar. Hech qanday izoh yozma.
"""
    raw = chat(prompt, temperature=0.6)
    raw = _strip_json_fence(raw)
    data = json.loads(raw)

    # Himoya: uzunligi kam/ko'p bo'lsa, kesib olamiz
    if isinstance(data, list):
        data = data[:QUESTIONS_PER_TEST]
    else:
        raise ValueError("Test JSON array bo'lishi kerak")
    return data


# ===================== PPTX BUILDER =====================
def make_pptx(slide_plan: dict, out_path: str):
    prs = Presentation()

    # Title slide
    title_layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(title_layout)
    s.shapes.title.text = slide_plan.get("title", "Prezentatsiya")
    if s.placeholders and len(s.placeholders) > 1:
        s.placeholders[1].text = "Talabalar uchun o‘quv materiali"

    # Content slides
    content_layout = prs.slide_layouts[1]  # Title and Content
    for item in slide_plan.get("slides", []):
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = item.get("title", "Slayd")

        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()

        bullets = item.get("bullets", [])
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = str(b)
            p.level = 0

        notes = item.get("notes", "")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    prs.save(out_path)


# ===================== TELEGRAM UI =====================
def main_menu() -> InlineKeyboardMarkup:
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("🎞 PPTX tayyorlash", callback_data="ASK_PPT_TOPIC")],
        [InlineKeyboardButton("📝 Test (25 ta / 25 daqiqa)", callback_data="ASK_TEST_TOPIC")],
    ])


def test_ready_menu(test_id: str) -> InlineKeyboardMarkup:
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("▶️ Testni boshlash", callback_data=f"START_TEST|{test_id}")],
        [InlineKeyboardButton("⬅️ Menu", callback_data="BACK_MENU")],
    ])


# ===================== HANDLERS =====================
async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(
        "Salom! Men Talaba AI bot.\n\n"
        "✅ Hamma natijalar PPTX bo‘ladi.\n"
        "✅ Test: 25 ta savol (Quiz) + 25 daqiqa.\n\n"
        "Boshlash: /menu"
    )


async def menu(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text("Kerakli bo‘limni tanlang:", reply_markup=main_menu())


async def ppt_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    topic = " ".join(context.args).strip()
    if not topic:
        await update.message.reply_text("Misol: /ppt Sun'iy intellekt ta'limda")
        return
    await create_and_send_pptx(update, context, topic)


async def create_and_send_pptx(update: Update, context: ContextTypes.DEFAULT_TYPE, topic: str):
    await update.message.reply_text("PPTX tayyorlanyapti...")

    try:
        plan = generate_slide_plan(topic)
        filename = f"presentation_{uuid.uuid4().hex[:8]}.pptx"
        out_path = os.path.join(".", filename)
        make_pptx(plan, out_path)

        await update.message.reply_document(
            document=open(out_path, "rb"),
            filename=filename,
            caption=f"✅ PPTX tayyor: {plan.get('title','Prezentatsiya')}",
        )
    except Exception as e:
        await update.message.reply_text(f"Xatolik: {e}")


async def on_button(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    await q.answer()

    if q.data == "ASK_PPT_TOPIC":
        context.user_data["mode"] = "ppt"
        await q.message.reply_text("Mavzuni yozing (masalan: 'Kiberxavfsizlik asoslari'):")
        return

    if q.data == "ASK_TEST_TOPIC":
        context.user_data["mode"] = "test"
        await q.message.reply_text("Test mavzusini yozing (masalan: 'SQL SELECT so‘rovlari'):")
        return

    if q.data == "BACK_MENU":
        await q.message.reply_text("Kerakli bo‘limni tanlang:", reply_markup=main_menu())
        return

    if q.data.startswith("START_TEST|"):
        _, test_id = q.data.split("|", 1)
        await start_test_run(update, context, test_id=test_id)
        return


async def handle_text(update: Update, context: ContextTypes.DEFAULT_TYPE):
    text = (update.message.text or "").strip()
    mode = context.user_data.get("mode")

    if mode == "ppt":
        context.user_data["mode"] = None
        await create_and_send_pptx(update, context, text)
        return

    if mode == "test":
        context.user_data["mode"] = None

        chat_id = update.message.chat_id
        user_id = update.message.from_user.id

        sessions = get_sessions(context.application)
        test_id = uuid.uuid4().hex[:10]

        sessions[(chat_id, user_id)] = {
            "test_id": test_id,
            "topic": text,
            "active": False,
            "started_at": None,
            "correct_count": 0,
            "answered_polls": set(),
            "poll_map": {},  # poll_id -> correct_idx
            "timer_job": None,
        }

        await update.message.reply_text(
            f"📝 Test mavzusi: {text}\n"
            f"— Savollar: {QUESTIONS_PER_TEST} ta\n"
            f"— Vaqt: 25 daqiqa\n\n"
            "Boshlash tugmasini bosing:",
            reply_markup=test_ready_menu(test_id),
        )
        return

    await update.message.reply_text("Buyruqlar: /menu yoki /ppt <mavzu>")


async def start_test_run(update: Update, context: ContextTypes.DEFAULT_TYPE, test_id: str):
    q = update.callback_query
    chat_id = q.message.chat_id
    user_id = q.from_user.id

    sessions = get_sessions(context.application)
    skey = (chat_id, user_id)
    sess = sessions.get(skey)

    if not sess or sess.get("test_id") != test_id:
        await q.message.reply_text("Bu test sessiyasi eskirib qolgan. /menu dan qaytadan boshlang.")
        return

    if sess.get("active"):
        await q.message.reply_text("Test allaqachon boshlangan.")
        return

    topic = sess.get("topic", "Umumiy")
    sess["active"] = True
    sess["started_at"] = _now()
    sess["correct_count"] = 0
    sess["answered_polls"] = set()
    sess["poll_map"] = {}

    # Eski timer bo'lsa olib tashlaymiz
    old_job = sess.get("timer_job")
    if old_job:
        try:
            old_job.schedule_removal()
        except Exception:
            pass

    await q.message.reply_text(
        f"⏱ Test boshlandi!\nMavzu: {topic}\nVaqt: 25 daqiqa\n\nSavollar ketma-ket chiqadi."
    )

    # Savollarni yaratamiz
    try:
        questions = generate_test_questions(topic)
    except Exception as e:
        sess["active"] = False
        await q.message.reply_text(f"Test savollarini yaratishda xatolik: {e}")
        return

    # Timer (25 daqiqadan keyin natija)
    job = context.job_queue.run_once(
        test_time_up,
        when=TEST_DURATION_SECONDS,
        data={"chat_id": chat_id, "user_id": user_id, "test_id": test_id},
        name=f"timer_{chat_id}_{user_id}_{test_id}",
    )
    sess["timer_job"] = job

    # Quiz poll yuboramiz
    letter_to_idx = {"A": 0, "B": 1, "C": 2, "D": 3}

    for i, item in enumerate(questions, start=1):
        qtext = (item.get("q") or "").strip()
        options = item.get("options") or []
        ans_letter = (item.get("answer") or "A").strip().upper()
        correct_idx = letter_to_idx.get(ans_letter, 0)

        clean_opts = []
        for opt in options[:4]:
            opt = str(opt).strip()
            # "A) text" -> "text"
            if ") " in opt:
                clean_opts.append(opt.split(") ", 1)[1].strip())
            else:
                clean_opts.append(opt)

        # Himoya: 4 ta variant bo'lmasa, to'ldirib qo'yamiz
        while len(clean_opts) < 4:
            clean_opts.append("—")

        poll_msg = await context.bot.send_poll(
            chat_id=chat_id,
            question=f"{i}) {qtext}",
            options=clean_opts[:4],
            type="quiz",
            correct_option_id=correct_idx,
            is_anonymous=False,
        )

        sess["poll_map"][poll_msg.poll.id] = correct_idx

    await q.message.reply_text("✅ 25 ta savol yuborildi. 25 daqiqadan so‘ng natija chiqadi.")


async def on_poll_answer(update: Update, context: ContextTypes.DEFAULT_TYPE):
    pa = update.poll_answer
    user_id = pa.user.id
    poll_id = pa.poll_id
    selected = pa.option_ids[0] if pa.option_ids else None
    if selected is None:
        return

    sessions = get_sessions(context.application)

    # Qaysi chat ekanini poll_answer’dan bilib bo'lmaydi.
    # Shuning uchun: userning barcha sessionlari ichidan poll_id mos kelganini topamiz.
    for (chat_id, uid), sess in sessions.items():
        if uid != user_id:
            continue
        if not sess.get("active"):
            continue

        poll_map = sess.get("poll_map", {})
        if poll_id not in poll_map:
            continue

        # Bir pollga 1 marta hisoblaymiz
        answered = sess.get("answered_polls", set())
        if poll_id in answered:
            return
        answered.add(poll_id)
        sess["answered_polls"] = answered

        correct_idx = poll_map[poll_id]
        if selected == correct_idx:
            sess["correct_count"] = sess.get("correct_count", 0) + 1
        return  # topdik, chiqamiz


async def test_time_up(context: ContextTypes.DEFAULT_TYPE):
    job = context.job
    chat_id = job.data["chat_id"]
    user_id = job.data["user_id"]
    test_id = job.data["test_id"]

    sessions = get_sessions(context.application)
    sess = sessions.get((chat_id, user_id))

    # Session bo'lmasa yoki boshqa test bo'lsa:
    if not sess or sess.get("test_id") != test_id:
        return

    # Testni yopamiz
    sess["active"] = False

    correct = int(sess.get("correct_count", 0))
    total = QUESTIONS_PER_TEST

    # Foiz
    percent = (correct / total) * 100 if total else 0

    if percent == 100:
        msg = f"🏆 OFARIN! 100%!\nTo‘g‘ri: {correct}/{total}"
    elif percent >= 90:
        msg = f"🔥 Juda yaxshi: {percent:.0f}%\nTo‘g‘ri: {correct}/{total}"
    elif percent >= 80:
        msg = f"✅ Yaxshi: {percent:.0f}%\nTo‘g‘ri: {correct}/{total}"
    else:
        msg = f"📌 Natija: {percent:.0f}%\nTo‘g‘ri: {correct}/{total}\nYana mashq qiling 🙂"

    await context.bot.send_message(chat_id=chat_id, text="⛔️ 25 daqiqa tugadi! Test yakunlandi.")
    await context.bot.send_message(chat_id=chat_id, text=msg)


def main():
    app = Application.builder().token(8261872993:AAGn79cuN2qAXAEeWXpl2txUqjzew61YrTg).build()

    # Commands
    app.add_handler(CommandHandler("start", start))
    app.add_handler(CommandHandler("menu", menu))
    app.add_handler(CommandHandler("ppt", ppt_command))

    # Buttons
    app.add_handler(CallbackQueryHandler(on_button))

    # Text input
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_text))

    # Quiz poll answers
    app.add_handler(PollAnswerHandler(on_poll_answer))

    app.run_polling()


if __name__ == "__main__":
    main()
